"""Generate downloadable PowerPoint decks from ACTION text."""

from __future__ import annotations

import io
import re
from typing import Any, Dict, List

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.media.copy_extract import extract_flyer_copy, extract_slide_bullets
from app.media.storage import GeneratedStorage


class PresentationBuilder:
    def __init__(self) -> None:
        self._storage = GeneratedStorage()

    async def build(
        self,
        *,
        user_id: str,
        assistant_text: str,
        title: str = "Blue Prince21 McKinzy Pitch",
    ) -> Dict[str, Any]:
        copy = extract_flyer_copy(assistant_text)
        bullets = extract_slide_bullets(assistant_text, limit=12)
        sections = self._section_slides(assistant_text)

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        # Title slide
        self._add_title_slide(
            prs,
            copy.get("headline") or title,
            copy.get("subhead") or copy.get("tagline") or "",
        )

        # Agenda / bullets
        if bullets:
            self._add_bullets_slide(prs, "Highlights", bullets[:6])
        if len(bullets) > 6:
            self._add_bullets_slide(prs, "Next steps", bullets[6:12])

        for sec_title, sec_bullets in sections[:6]:
            self._add_bullets_slide(prs, sec_title, sec_bullets[:6])

        # Closing
        self._add_title_slide(
            prs,
            copy.get("cta") or "Let's build together",
            "anthonywarrenmckinzy.com · Drink it. Trade it. Own it.",
        )

        buf = io.BytesIO()
        prs.save(buf)
        data = buf.getvalue()
        path, url = await self._storage.upload_bytes(
            user_id=user_id,
            data=data,
            ext="pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            folder="decks",
        )
        asset = {
            "kind": "other",
            "title": f"{title}.pptx",
            "storage_path": path,
            "public_url": url,
            "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "byte_size": len(data),
            "meta": {"format": "pptx", "print_ready": False, "download": True},
        }
        return {"assets": [asset], "primary_url": url}

    @staticmethod
    def _section_slides(text: str) -> List[tuple[str, List[str]]]:
        parts = re.split(r"\n(?=#{1,3}\s+|\*\*[A-Z][^*]{3,40}\*\*)", text or "")
        out: List[tuple[str, List[str]]] = []
        for part in parts:
            title_m = re.search(r"^(?:#{1,3}\s+|\*\*)(.+?)(?:\*\*)?\s*$", part.strip(), re.M)
            title = title_m.group(1).strip() if title_m else "Overview"
            bullets = extract_slide_bullets(part, limit=6)
            if not bullets:
                para = re.sub(r"[#*_`]", " ", part)
                para = re.sub(r"\s+", " ", para).strip()
                if len(para) > 60:
                    bullets = [para[:200]]
            if bullets:
                out.append((title[:60], bullets))
        return out

    @staticmethod
    def _style_shape_bg(shape, rgb=(11, 61, 46)) -> None:
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*rgb)

    def _add_title_slide(self, prs: Presentation, title: str, subtitle: str) -> None:
        layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(layout)
        # background via rectangle
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        self._style_shape_bg(shape)
        shape.line.fill.background()

        box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(2))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title[:120]
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(245, 240, 230)
        p.alignment = PP_ALIGN.LEFT

        sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.4), Inches(11.5), Inches(1.5))
        st = sub.text_frame
        st.word_wrap = True
        sp = st.paragraphs[0]
        sp.text = subtitle[:200]
        sp.font.size = Pt(20)
        sp.font.color.rgb = RGBColor(201, 162, 39)

    def _add_bullets_slide(
        self, prs: Presentation, title: str, bullets: List[str]
    ) -> None:
        layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(layout)
        bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), prs.slide_width, prs.slide_height
        )
        self._style_shape_bg(bg, (26, 26, 26))
        bg.line.fill.background()

        tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.5), Inches(1))
        tp = tbox.text_frame.paragraphs[0]
        tp.text = title[:80]
        tp.font.size = Pt(28)
        tp.font.bold = True
        tp.font.color.rgb = RGBColor(201, 162, 39)

        body = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(5.2))
        tf = body.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {b}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(245, 240, 230)
            p.space_after = Pt(10)
